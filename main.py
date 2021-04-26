# This is a sample Python script.

# Press Shift+F10 to execute it or replace it with your code.
# Press Double Shift to search everywhere for classes, files, tool windows, actions, and settings.

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.dml.fill import FillFormat
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# main
prs = Presentation('testing.pptx')


# Child and mother slide function, first slide on presentation
def first_slide():
    # First slide
    slide1 = prs.slide_masters[6].slide_layouts[3]
    # Add child and mom slide
    slide1 = prs.slides.add_slide(slide1)
    first_slide = slide1.placeholders[0]
    first_slide.left = Inches(0.5)
    first_slide.top = Inches(3)
    first_slide.width = Inches(8)
    first_slide.height = Inches(1)

    # Set text box for slide 1
    for shape in slide1.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Sainsbury Wan Service Review\n" \
                   "April 2021\n" \
                   "Simon Hardaker\n" \
                   "Reporting Period:  March 2021"
        font = run.font
        font.name = 'Arial'
        font.size = Pt(32)
        font.bold = False
        font.italic = None  # cause value to be inherited from theme


def second_slide():
    # Index slide layout
    slide2 = prs.slide_masters[6].slide_layouts[18]

    # index slide
    slide2 = prs.slides.add_slide(slide2)
    # title format
    # Slide 2
    index_title = slide2.placeholders[10]
    index_title.text = "Index"
    index_title.left = Inches(0.5)
    index_title.top = Inches(0.3)
    index_title.width = Inches(8)
    index_title.height = Inches(1)

    # textbox format
    bullet_point_box = slide2.shapes
    bullet_point1 = bullet_point_box.placeholders[14]
    bullet_point1.left = Inches(0.5)
    bullet_point1.top = Inches(2)
    bullet_point1.width = Inches(8)
    bullet_point1.height = Inches(1)
    bullet_point1.text = "Executive Summary\n" \
                         "Points for Discussion\n" \
                         "Site Breakdown\n" \
                         "Interface Availability \n" \
                         "Bandwidth Utilization\n" \
                         "Incident Management\n" \
                         "Change Management\n" \
                         "Problem Management\n" \
                         "Service Improvement Summary\n" \
                         "Emergency contacts\n" \
                         "End of contract information\n"


# Third slide with a table
def third_slide():
    # Table slide layout
    slide_table = prs.slide_masters[6].slide_layouts[18]

    # table slide
    slide3 = prs.slides.add_slide(slide_table)

    # Set title for slide 3
    slide3_title = slide3.placeholders[10]
    slide3_title.text = 'Executive Summary'
    slide3_title.left = Inches(0.5)
    slide3_title.top = Inches(0.3)
    slide3_title.width = Inches(8)
    slide3_title.height = Inches(1)

    shapes = slide3.shapes

    # ---add table to slide---

    rows = 3
    cols = 2
    left = Inches(0.3)
    top = Inches(1.5)
    width = Inches(6.0)
    height = Inches(0.8)

    table = shapes.add_table(rows, cols, left, top, width, height).table

    # set column widths
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(10.0)

    # set row heights
    table.rows[0].height = Inches(0)
    table.rows[1].height = Inches(2)
    table.rows[2].height = Inches(1.5)

    # colours of cell titles
    table.cell(0, 0).fill.solid()
    table.cell(0, 0).fill.fore_color.rgb = RGBColor(25, 34, 79)
    table.cell(0, 1).fill.solid()
    table.cell(0, 1).fill.fore_color.rgb = RGBColor(25, 34, 79)

    # colours of cell columns
    table.cell(1, 0).fill.solid()
    table.cell(1, 0).fill.fore_color.rgb = RGBColor(242, 242, 242)

    table.cell(1, 1).fill.solid()
    table.cell(1, 1).fill.fore_color.rgb = RGBColor(242, 242, 242)

    table.cell(2, 0).fill.solid()
    table.cell(2, 0).fill.fore_color.rgb = RGBColor(242, 242, 242)

    table.cell(2, 1).fill.solid()
    table.cell(2, 1).fill.fore_color.rgb = RGBColor(242, 242, 242)

    # write column headings
    table.cell(0, 0).text = "Management Summary"

    table.cell(0, 1).text = 'Description'

    # write body cells
    table.cell(1, 0).text = "Highlights"
    table.cell(1, 1).text = "- Availability – Green across the board.\n" \
                            "- Incident volumes reduce further.This trend has continued in Feb\n" \
                            "- Repeat site issues extremely low.\n" \
                            "- All incident priorities achieved SLA.\n" \
                            "- MTTR improved across all INC priorities.\n"

    table.cell(2, 0).text = "Lowlights"
    table.cell(2, 1).text = "- Emerson’s Green: Appearing on both the utilisation observations and repeat INC slide.\n" \
                            "- Fault not found, logged in error etc make up over a third of the overall incident ticket " \
                            "volume. \n" \
                            "- New order quotes: SLA missed this month."


def fourth_slide():
    slide = prs.slide_masters[6].slide_layouts[18]
    slide = prs.slides.add_slide(slide)
    # define chart data ---------------------
    chart_data = CategoryChartData()
    chart_data.categories = ['East', 'West', 'Midwest']
    chart_data.add_series('Series 1', (19.2, 21.4, 16.7))

    # add chart to slide --------------------
    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    )


# # Slide placeholders
# for shape in slide1.placeholders:
#     print('%d %s' % (shape.placeholder_format.idx, shape.shape_type))

# print("for placeholders format and type")
# for shape in slide3.shapes:
#     if shape.is_placeholder:
#         phf = shape.placeholder_format
#         print('%d, %s' % (phf.idx, phf.type))
first_slide()
second_slide()
third_slide()
fourth_slide()

prs.save('testing.pptx')
